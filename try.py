import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
import openai
import os
from azure.storage.blob import BlobServiceClient, ContainerClient
from io import BytesIO

# Set your OpenAI API key
openai.api_type = "azure"
openai.api_version = "2023-07-01-preview"
openai.api_base = "https://hsbcoai.openai.azure.com/"
openai.api_key = "324de9c5fb0444278d98f92dbee4ee2a"

# Azure Blob Storage connection string
azure_storage_connection_string = "DefaultEndpointsProtocol=https;AccountName=demodocstoragetranslator;AccountKey=cr4tK/oThygMWmBm3YvfuSxIwj8aMiMPiF+jz27xwpQTPDJxmi8LZe7UxrqsgXcd8Ggtnl/6KrdF+AStMqVAEQ==;EndpointSuffix=core.windows.net"
container_name = "ppt"  # Replace with your container name

def generate_content(prompt):
    # Use OpenAI to generate content based on the prompt
    response = openai.Completion.create(
        engine="gpt-35-turbo",
        prompt=prompt,
        max_tokens=140,  # Adjust as needed
        temperature=0.2,  # Adjust for creativity
    )
    return response.choices[0].text

def create_presentation(topic, selected_titles):
    # Download the template from Azure Blob Storage
    try:
        blob_service_client = BlobServiceClient.from_connection_string(azure_storage_connection_string)
        container_client = blob_service_client.get_container_client(container_name)
        blob_client = container_client.get_blob_client("template.pptx")

        template_presentation_data = blob_client.download_blob()
        template_presentation_path = "template.pptx"

        with open(template_presentation_path, "wb") as data:
            data.write(template_presentation_data.readall())

    except Exception as e:
        return str(e)

    # Create the presentation based on the template
    prs = Presentation(template_presentation_path)

    slide_titles = selected_titles
    slide_layouts = prs.slide_layouts

    for title in slide_titles:
        slide_layout = slide_layouts[1]  # Use Title Slide layout (index 1) for slides
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Change the font type of titles to something else and make them bold
        title_text_frame = title_placeholder.text_frame
        title_text_frame.paragraphs[0].font.name = 'Palace Script MT'  # Change to your desired font
        title_text_frame.paragraphs[0].font.bold = True

        # Generate content for the content placeholder using OpenAI
        content_prompt = f"Write 4 key points about the '{topic}' for '{title}' slide. The points should only be related to the  '{topic}' and no irrelevant information should be added:"
        content = generate_content(content_prompt)

        # Add bullet points with a maximum of 4 points per slide
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        paragraphs = content.split('\n')
        max_points = 4
        points_added = 0

        for paragraph in paragraphs:
            if paragraph.strip():
                p = text_frame.add_paragraph()
                p.text = paragraph
                p.space_after = Pt(12)
                p.space_before = Pt(0)
                p.level = 0

                # Set the font size to 18 Pt and change to Times New Roman
                for run in p.runs:
                    font = run.font
                    font.size = Pt(10)
                    font.name = 'Times New Roman'

                points_added += 1

                if points_added >= max_points:
                    break

    # Add the HSBC logo to all slides
    logo_path = 'hsbc.png'
    left = Inches(8.2)  # Adjust the left position as needed
    top = Inches(0.5)  # Adjust the top position as needed
    for slide in prs.slides:
        pic = slide.shapes.add_picture(logo_path, left, top, height=Inches(1))  # Adjust the height as needed

    # Replace the topic on the first slide
    if prs.slides:
        first_slide = prs.slides[0]
        title_placeholder = first_slide.shapes.title
        title_placeholder.text = topic
        # Set the title text alignment to center
        title_text_frame = title_placeholder.text_frame
        title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save the generated presentation
    generated_presentation_path = f"{topic}_presentation.pptx"
    prs.save(generated_presentation_path)

    # Delete the template file
    os.remove(template_presentation_path)

    # Upload the generated presentation to Azure Blob Storage
    blob_client = container_client.get_blob_client(generated_presentation_path)
    with open(generated_presentation_path, "rb") as data:
        blob_client.upload_blob(data, overwrite=True)

    return generated_presentation_path

def delete_slide(presentation, slide_index):
    slides = presentation.slides
    if 0 <= slide_index < len(slides):
        slides._sldIdLst.remove(slides._sldIdLst[slide_index])

def download_presentation_from_azure_blob(blob_name, download_name):
    try:
        blob_service_client = BlobServiceClient.from_connection_string(azure_storage_connection_string)
        container_client = blob_service_client.get_container_client(container_name)
        blob_client = container_client.get_blob_client(blob_name)

        with open(download_name, "wb") as download_data:
            download_data.write(blob_client.download_blob().readall())

        return True
    except Exception as e:
        return str(e)

def main():
    st.title("PowerPoint Presentation Generator")

    # Get user input for the presentation topic
    topic = st.text_input("Enter the topic for the presentation:")

    # Create a select box for choosing slide titles
    available_slide_titles = ["Introduction", "Main Content", "Advantages", "Disadvantages", "Conclusion"]
    selected_titles = st.multiselect("Select Slide Titles", available_slide_titles)

    if st.button("Generate Presentation"):
        if topic and selected_titles:
            st.text("Generating the presentation. Please wait...")

            # Create the presentation
            presentation_path = create_presentation(topic, selected_titles)
            st.success(f"Presentation generated and saved in Azure Blob Storage as '{presentation_path}'")

    if st.button("Download Presentation"):
        if topic:
            download_name = f"{topic}_presentation.pptx"
            if download_presentation_from_azure_blob(download_name, download_name):
                st.success(f"Presentation downloaded as '{download_name}'")
                st.text("Deleting the second slide...")
                presentation = Presentation(download_name)
                delete_slide(presentation, 1)  # Delete the second slide (index 1)
                presentation.save(download_name)
                st.success("Second slide deleted and presentation saved.")
            else:
                st.error("Error downloading the presentation")

if __name__ == "__main__":
    main()




